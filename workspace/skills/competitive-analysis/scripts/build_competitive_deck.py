#!/usr/bin/env python3
"""
生成竞争格局分析的基础 PPTX 模型。
优先从本地 Excel 提取数据，缺失数据留空并标记 TODO 供 Web Search 填补。
遵循 SKILL.md 中的字体和颜色规范。
"""
import argparse
import re
import zipfile
import tempfile
import os
from pathlib import Path
from typing import Any, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError as exc:
    raise ImportError("python-pptx required: pip install python-pptx") from exc

try:
    import openpyxl
except ImportError as exc:
    raise ImportError("openpyxl required: pip install openpyxl") from exc

# ==================== Data Extraction Logic ====================
def safe_divide(n, d): return n / d if d != 0 else 0.0

def _ensure_shared_strings(file_path: Path) -> Tuple[Path, bool]:
    file_path = Path(file_path)
    try:
        with zipfile.ZipFile(str(file_path), 'r') as zf:
            if 'xl/sharedStrings.xml' in zf.namelist(): return file_path, False
            ct = zf.read('[Content_Types].xml').decode('utf-8', errors='ignore')
            if 'sharedStrings' not in ct: return file_path, False
    except: return file_path, False

    fd, tmp = tempfile.mkstemp(suffix='.xlsx', dir=str(file_path.parent))
    os.close(fd)
    try:
        with zipfile.ZipFile(str(file_path), 'r') as src, zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as dst:
            for item in src.namelist(): dst.writestr(item, src.read(item))
            dst.writestr('xl/sharedStrings.xml', '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="0" uniqueCount="0"/>')
        return Path(tmp), True
    except: return file_path, False

def _read_excel_map(file_path: Path) -> dict:
    fp, is_tmp = _ensure_shared_strings(file_path)
    try:
        wb = openpyxl.load_workbook(str(fp), read_only=True, data_only=True)
        sheet = wb.active
        data = {}
        header_row = 1
        for r in range(1, 5):
            if "FY" in " ".join(str(sheet.cell(r, c).value) for c in range(1, sheet.max_column+1) if sheet.cell(r,c).value):
                header_row = r; break

        periods = [str(sheet.cell(header_row, c).value) for c in range(3, sheet.max_column + 1)]
        for r in range(header_row + 1, sheet.max_row + 1):
            ind = sheet.cell(r, 1).value
            if not ind: continue
            ind = str(ind).strip()
            vals = {}
            for c, p in enumerate(periods):
                v = sheet.cell(r, c + 3).value
                num = 0.0
                if isinstance(v, (int, float)): num = float(v)
                elif isinstance(v, str) and v:
                    cl = v.replace(",", "").replace("-", "").strip()
                    if cl.replace(".", "", 1).isdigit(): num = float(cl) * (-1 if v.strip().startswith("-") else 1)
                vals[p] = num
            data[ind] = vals
        wb.close()
        return data
    finally:
        if is_temp and fp.exists(): os.remove(str(fp))

def find_local_file(workspace: Path, ticker: str, suffix: str) -> Path:
    pattern = re.compile(rf'^.*_{re.escape(ticker)}_{re.escape(suffix)}_.*\.(xlsx|xls)$', re.IGNORECASE)
    files = [f for f in (workspace / "excels").iterdir() if f.is_file() and pattern.match(f.name)]
    return sorted(files)[-1] if files else None

def extract_company_data(ticker: str, workspace: Path) -> dict:
    data = {"ticker": ticker, "revenue": None, "net_income": None, "source": "Web Search Required"}
    inc_file = find_local_file(workspace, ticker, "income")

    if inc_file:
        inc_data = _read_excel_map(inc_file)
        fy_cols = sorted([c for c in inc_data.get("总收入", {}).keys() if "FY" in c], reverse=True)
        if fy_cols:
            latest = fy_cols[0]
            data["revenue"] = inc_data.get("总收入", {}).get(latest)
            data["net_income"] = inc_data.get("净利润", {}).get(latest)
            data["source"] = f"Local file: {inc_file.name}"
    return data

# ==================== PPTX Builder Class ====================
class CompetitiveDeckBuilder:
    def __init__(self, target: str, competitors: list, workspace: Path):
        self.target = target
        self.competitors = competitors
        self.workspace = workspace
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def build(self, output_path) -> Path:
        target_data = extract_company_data(self.target, self.workspace)
        comp_data_list = [extract_company_data(t, self.workspace) for t in self.competitors]

        self._build_title_slide()
        self._build_target_profile(target_data)
        self._build_comparative_table(target_data, comp_data_list)

        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        print(f"Competitive Deck saved to {path}")
        return path

    def _add_textbox(self, slide, left, top, width, height, text, size=14, bold=False, color="000000", align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor.from_string(color)
        p.alignment = align
        return txBox

    def _build_title_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Blank
        # Title
        self._add_textbox(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.5),
                          f"{self.target} Competitive Landscape", size=32, bold=True, color="1F4E79", align=PP_ALIGN.CENTER)
        # Subtitle
        comps_str = ", ".join(self.competitors)
        self._add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(1.0),
                          f"Comparing against: {comps_str}", size=18, color="595959", align=PP_ALIGN.CENTER)

    def _build_target_profile(self, data: dict):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # Title
        self._add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
                          "Target Company Profile", size=28, bold=True, color="1F4E79")

        # Table
        rows, cols = 4, 2
        left, top = Inches(1), Inches(1.5)
        width, height = Inches(6), Inches(3)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Headers
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        for c in range(2):
            cell = table.cell(0, c)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string("F2F2F2")

        # Data
        metrics = [
            ("Ticker", data["ticker"]),
            ("Revenue", f"{data['revenue']:,}" if data['revenue'] else "TODO: Web Search"),
            ("Net Income", f"{data['net_income']:,}" if data['net_income'] else "TODO: Web Search")
        ]

        for i, (m, v) in enumerate(metrics):
            table.cell(i+1, 0).text = m
            table.cell(i+1, 1).text = str(v)
            table.cell(i+1, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string("0000FF" if "TODO" in str(v) else "000000")

    def _build_comparative_table(self, target_data: dict, comp_data_list: list):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # Title
        self._add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
                          "Comparative Financial Analysis", size=28, bold=True, color="1F4E79")

        all_comps = [target_data] + comp_data_list
        rows, cols = 4, len(all_comps) + 1
        left, top = Inches(0.5), Inches(1.5)
        width, height = Inches(12.3), Inches(3.0)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Headers
        table.cell(0, 0).text = "Metric"
        for i, c in enumerate(all_comps):
            table.cell(0, i+1).text = c["ticker"]

        for c in range(cols):
            cell = table.cell(0, c)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string("F2F2F2")

        # Data
        metrics = ["Revenue", "Net Income", "Source"]
        for i, m in enumerate(metrics):
            table.cell(i+1, 0).text = m
            for j, c in enumerate(all_comps):
                val = ""
                if m == "Revenue": val = f"{c['revenue']:,}" if c['revenue'] else "TODO"
                elif m == "Net Income": val = f"{c['net_income']:,}" if c['net_income'] else "TODO"
                else: val = c['source']

                table.cell(i+1, j+1).text = val
                if "TODO" in val or "Web Search" in val:
                    table.cell(i+1, j+1).text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string("FF0000")
                else:
                    table.cell(i+1, j+1).text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string("0000FF")

def main() -> None:
    parser = argparse.ArgumentParser(description="Build Competitive Landscape PPTX")
    parser.add_argument("--target", required=True, help="Target company ticker")
    parser.add_argument("--competitors", required=True, help="Comma-separated competitor tickers")
    parser.add_argument("--workspace", required=True, help="Path to workspace directory")
    parser.add_argument("--output", required=True, help="Path to output pptx file")
    args = parser.parse_args()

    builder = CompetitiveDeckBuilder(args.target, args.competitors.split(','), Path(args.workspace))
    builder.build(args.output)

if __name__ == "__main__":
    main()
